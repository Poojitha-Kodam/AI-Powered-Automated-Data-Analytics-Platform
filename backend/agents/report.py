import os
import pandas as pd
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
import docx
from pptx import Presentation
from pptx.util import Inches, Pt
from config import settings

class ReportAgent:
    def __init__(self):
        pass

    def generate_pdf_report(self, filepath: str, title: str, summary: str, columns_list: list, clean_history: list, insights: list, forecast: dict = None) -> str:
        """
        Generate a professional PDF report containing the executive summary, cleaning details, and insights.
        """
        doc = SimpleDocTemplate(filepath, pagesize=letter)
        styles = getSampleStyleSheet()
        
        # Custom styles for dark theme look or sleek corporate design
        title_style = ParagraphStyle(
            'TitleStyle',
            parent=styles['Heading1'],
            fontSize=24,
            leading=28,
            textColor=colors.HexColor('#1e1b4b'), # Deep Indigo
            spaceAfter=15
        )
        subtitle_style = ParagraphStyle(
            'SubtitleStyle',
            parent=styles['Normal'],
            fontSize=10,
            textColor=colors.HexColor('#475569'),
            spaceAfter=20
        )
        h2_style = ParagraphStyle(
            'H2Style',
            parent=styles['Heading2'],
            fontSize=16,
            leading=20,
            textColor=colors.HexColor('#4f46e5'), # Indigo
            spaceBefore=15,
            spaceAfter=10
        )
        body_style = ParagraphStyle(
            'BodyStyle',
            parent=styles['Normal'],
            fontSize=10,
            leading=14,
            textColor=colors.HexColor('#1e293b'),
            spaceAfter=8
        )
        bullet_style = ParagraphStyle(
            'BulletStyle',
            parent=styles['Normal'],
            fontSize=10,
            leading=14,
            textColor=colors.HexColor('#1e293b'),
            leftIndent=15,
            firstLineIndent=-10,
            spaceAfter=5
        )

        story = []
        
        # Title
        story.append(Paragraph(title, title_style))
        story.append(Paragraph(f"Generated on: {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S')}", subtitle_style))
        story.append(Spacer(1, 10))

        # Executive Summary
        story.append(Paragraph("Executive Summary", h2_style))
        story.append(Paragraph(summary, body_style))
        story.append(Spacer(1, 10))

        # Dataset Structure
        story.append(Paragraph("Dataset Schema", h2_style))
        columns_text = ", ".join(columns_list)
        story.append(Paragraph(f"<b>Columns Analyzed:</b> {columns_text}", body_style))
        story.append(Spacer(1, 10))

        # Cleaning Actions Taken
        story.append(Paragraph("Data Cleaning Summary", h2_style))
        if clean_history:
            for item in clean_history:
                story.append(Paragraph(f"• {item}", bullet_style))
        else:
            story.append(Paragraph("No cleaning anomalies detected; dataset was ready for analysis.", body_style))
        story.append(Spacer(1, 10))

        # Insights
        story.append(Paragraph("Automated Data Insights", h2_style))
        for ins in insights:
            story.append(Paragraph(f"• {ins}", bullet_style))
        story.append(Spacer(1, 10))

        # Forecast
        if forecast and not forecast.get("error"):
            story.append(Paragraph("Time Series Forecasting", h2_style))
            story.append(Paragraph(f"<b>Key Trend Recommendation:</b> {forecast.get('recommendation', '')}", body_style))
            
            # Simple table of predictions
            table_data = [["Date", "Forecasted Value", "Lower Limit", "Upper Limit"]]
            f_dates = forecast.get("forecast_dates", [])
            f_vals = forecast.get("forecast_values", [])
            f_lows = forecast.get("lower_bound", [])
            f_highs = forecast.get("upper_bound", [])
            
            for d, v, l, h in zip(f_dates[:6], f_vals[:6], f_lows[:6], f_highs[:6]):
                table_data.append([d, f"{v:.2f}", f"{l:.2f}", f"{h:.2f}"])
                
            t = Table(table_data, colWidths=[120, 100, 100, 100])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#6366f1')),
                ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
                ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                ('BOTTOMPADDING', (0,0), (-1,0), 6),
                ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#f8fafc')),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
                ('FONTSIZE', (0,0), (-1,-1), 9),
            ]))
            story.append(t)

        doc.build(story)
        return filepath

    def generate_docx_report(self, filepath: str, title: str, summary: str, columns_list: list, clean_history: list, insights: list, forecast: dict = None) -> str:
        """
        Generate a Microsoft Word report.
        """
        doc = docx.Document()
        doc.add_heading(title, 0)
        
        doc.add_heading("Executive Summary", level=1)
        doc.add_paragraph(summary)

        doc.add_heading("Dataset Schema", level=1)
        doc.add_paragraph(f"Columns: {', '.join(columns_list)}")

        doc.add_heading("Data Cleaning Actions", level=1)
        if clean_history:
            for item in clean_history:
                doc.add_paragraph(item, style='List Bullet')
        else:
            doc.add_paragraph("No cleaning actions needed. Dataset was fully clean.")

        doc.add_heading("Automated Insights", level=1)
        for ins in insights:
            doc.add_paragraph(ins, style='List Bullet')

        if forecast and not forecast.get("error"):
            doc.add_heading("Future Forecasting & Trends", level=1)
            doc.add_paragraph(forecast.get("recommendation", ""))
            
            table = doc.add_table(rows=1, cols=4)
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = 'Date'
            hdr_cells[1].text = 'Forecasted Value'
            hdr_cells[2].text = 'Lower Limit'
            hdr_cells[3].text = 'Upper Limit'
            
            f_dates = forecast.get("forecast_dates", [])
            f_vals = forecast.get("forecast_values", [])
            f_lows = forecast.get("lower_bound", [])
            f_highs = forecast.get("upper_bound", [])
            
            for d, v, l, h in zip(f_dates[:12], f_vals[:12], f_lows[:12], f_highs[:12]):
                row_cells = table.add_row().cells
                row_cells[0].text = d
                row_cells[1].text = f"{v:.2f}"
                row_cells[2].text = f"{l:.2f}"
                row_cells[3].text = f"{h:.2f}"

        doc.save(filepath)
        return filepath

    def generate_pptx_presentation(self, filepath: str, title: str, summary: str, insights: list, forecast: dict = None) -> str:
        """
        Generate a PowerPoint Presentation deck.
        """
        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_box = slide.shapes.title
        subtitle_box = slide.placeholders[1]
        
        title_box.text = title
        subtitle_box.text = f"Automated AI Analytics Presentation\nCreated on: {pd.Timestamp.now().strftime('%Y-%m-%d')}"

        # Slide 2: Executive Summary
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        shapes.title.text = "Executive Summary"
        tf = shapes.placeholders[1].text_frame
        tf.text = summary

        # Slide 3: Insights
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        shapes.title.text = "Key Data Insights"
        tf = shapes.placeholders[1].text_frame
        for ins in insights[:5]: # top 5 insights
            p = tf.add_paragraph()
            p.text = ins
            p.level = 0

        # Slide 4: Forecasting
        if forecast and not forecast.get("error"):
            slide = prs.slides.add_slide(slide_layout)
            shapes = slide.shapes
            shapes.title.text = "Future Forecasting & Recommendations"
            tf = shapes.placeholders[1].text_frame
            tf.text = forecast.get("recommendation", "")
            
            p = tf.add_paragraph()
            p.text = f"Forecast Period: {len(forecast.get('forecast_dates', []))} steps ahead"
            p.level = 1
            
            p2 = tf.add_paragraph()
            avg_forecast = sum(forecast.get('forecast_values', [])) / len(forecast.get('forecast_values', []))
            p2.text = f"Average predicted value over future period: {avg_forecast:.2f}"
            p2.level = 1

        prs.save(filepath)
        return filepath
        

